#!/usr/bin/env python3
"""
assemble_figure.py — fill an SVG template from a manifest, normalize fonts,
inject CC-BY attribution, and export SVG / PDF / PPTX.

Workflow:
  template SVG  +  manifest.yaml  ->  filled SVG  ->  PDF / PPTX

  manifest.panels: { placeholder-id : element-id-or-relative-path }
  manifest.labels: { text-element-id : "string content" }

Dependencies:
  - lxml, pyyaml         (always)
  - Inkscape >= 1.2 CLI  (PDF / PNG export)
  - python-pptx          (PPTX export, optional)

Usage:
  python assemble_figure.py \
      --template templates/dark_proteome_4panel_template.svg \
      --manifest templates/dark_proteome_manifest.yaml \
      --output ./fig1.svg \
      --export-pdf --export-pptx
"""

from __future__ import annotations

import argparse
import copy
import json
import re
import subprocess
import sys
from pathlib import Path

import yaml
from lxml import etree


SVG_NS = {
    "svg": "http://www.w3.org/2000/svg",
    "xlink": "http://www.w3.org/1999/xlink",
}
SVG_TAG = "{http://www.w3.org/2000/svg}"
XLINK_HREF = "{http://www.w3.org/1999/xlink}href"


# ---------------------------------------------------------------------------
# Manifest + element resolution
# ---------------------------------------------------------------------------


def load_manifest(manifest_path: Path) -> dict:
    """Load the YAML manifest as a plain dict.

    yaml.safe_load returns None for an empty / comment-only file; coerce to
    an empty dict so downstream `manifest.get(...)` calls don't blow up with
    AttributeError.
    """
    return yaml.safe_load(manifest_path.read_text()) or {}


# Same shape that library_tools.load_unified_index() emits for legacy YAML rows.
_LEGACY_CATEGORIES = (
    "cells", "molecules", "organelles", "tissues",
    "organs", "equipment", "pathways", "arrows",
)


def _scan_legacy_yaml(value: str, library_root: Path) -> dict | None:
    """Walk per-category _index.yaml files for a `manual-...` style ID.

    library_tools.search emits IDs like `manual-cells-treg-flat-v1` for rows
    registered through the legacy YAML path; without this fallback, those IDs
    can't be resolved here, even though the same ID renders in `search`.
    """
    if not value.startswith("manual-"):
        return None
    for category in _LEGACY_CATEGORIES:
        yaml_index = library_root / category / "_index.yaml"
        if not yaml_index.exists():
            continue
        try:
            data = yaml.safe_load(yaml_index.read_text()) or {"elements": []}
        except Exception:
            continue
        for entry in data.get("elements", []):
            file_stem = (entry.get("file") or "").replace(".svg", "")
            fid = f"manual-{category}-{file_stem}"
            if fid == value:
                return {
                    "id": fid,
                    "_category": category,
                    "_file": entry.get("file", ""),
                    "license": entry.get("license", "research-use-only"),
                    "source_name": entry.get("provider", "manual"),
                    "attribution_required": False,
                }
    return None


def resolve_element_path(value: str, library_root: Path) -> Path:
    """Map a manifest value to an actual on-disk file.

    Disambiguation:
      * "/" in value or extension in {.svg, .png, .emf}  ->  treated as a
        path relative to `library_root`.
      * starts with "manual-"                             ->  legacy YAML lookup.
      * else                                              ->  unified-index ID
        lookup in library_root/library/index.json.
    """
    if "/" in value or value.lower().endswith((".svg", ".png", ".emf")):
        return library_root / value

    legacy = _scan_legacy_yaml(value, library_root)
    if legacy is not None:
        return library_root / legacy["_category"] / legacy["_file"]

    json_index = library_root / "library" / "index.json"
    if json_index.exists():
        try:
            records: list[dict] = json.loads(json_index.read_text())
            for r in records:
                if r.get("id") == value:
                    return library_root / "library" / r["file"]
        except Exception:
            pass
    return library_root / value


def lookup_attribution(value: str, library_root: Path) -> dict | None:
    """If `value` references a CC-BY element, return its record.

    Resolves both shapes:
      * unified-index ID  ->  match `r["id"] == value`
      * relative path     ->  match `<library_root>/library/<r["file"]>`
                              against the resolved manifest path; this covers
                              the case where a user references a Reactome /
                              Servier asset by path instead of ID.

    Walks the unified `library/index.json` and the legacy per-category YAMLs.
    Legacy YAML rows aren't currently flagged attribution_required, but the
    symmetric scan keeps the resolution model consistent if that changes.
    """
    json_index = library_root / "library" / "index.json"

    # ID branch.
    if not ("/" in value or value.lower().endswith((".svg", ".png", ".emf"))):
        if json_index.exists():
            try:
                for r in json.loads(json_index.read_text()):
                    if r.get("id") == value and r.get("attribution_required"):
                        return r
            except Exception:
                pass
        legacy = _scan_legacy_yaml(value, library_root)
        if legacy and legacy.get("attribution_required"):
            return legacy
        return None

    # Path branch: resolve and reverse-lookup against the unified index by
    # comparing against the on-disk path each record points at.
    target = (library_root / value).resolve()
    if json_index.exists():
        try:
            for r in json.loads(json_index.read_text()):
                if not r.get("attribution_required"):
                    continue
                rec_path = (library_root / "library" / r["file"]).resolve()
                if rec_path == target:
                    return r
        except Exception:
            pass
    return None


# ---------------------------------------------------------------------------
# Template filling
# ---------------------------------------------------------------------------


def _find_placeholders(root: etree._Element, panel_id: str) -> list[etree._Element]:
    """Return SVG nodes whose `data-placeholder` attr equals `panel_id`.

    Uses lxml's `findall` with attribute predicate to avoid dynamic xpath
    string injection from manifest keys.
    """
    return root.findall(
        f'.//svg:*[@data-placeholder="{panel_id}"]',
        namespaces=SVG_NS,
    ) if not _suspicious_id(panel_id) else []


def _suspicious_id(panel_id: str) -> bool:
    """Reject IDs containing chars that could break the xpath literal."""
    return any(c in panel_id for c in ('"', "\n", "\r"))


# Attributes that may carry a `url(#...)` reference or a bare `#id`.
_REF_ATTRS = (
    "href", "{http://www.w3.org/1999/xlink}href",
    "clip-path", "mask", "filter",
    "fill", "stroke",
    "marker-start", "marker-mid", "marker-end",
)


def _namespace_subtree_ids(wrapper: etree._Element, prefix: str) -> None:
    """Prefix every `id` and same-document reference inside `wrapper`.

    When the same source SVG is dropped into multiple placeholders, deep-copying
    its `<defs>` / `<clipPath>` / `<mask>` / `<linearGradient>` etc. produces
    duplicate `id`s and `url(#xxx)` references that bind to whichever copy XML
    visits first — wrong renders, broken selectors. We rewrite both sides.

    The wrapper element itself is skipped: its id was set to `prefix` by the
    caller, and re-prefixing would produce e.g. `panel-a-1-panel-a-1`, breaking
    every downstream `dom_set` selector that targets panels by their expected
    id.
    """
    # 1) Collect ids and rewrite the id attribute itself.
    #    Skip the wrapper itself (its id is the caller-supplied prefix).
    id_map: dict[str, str] = {}
    for elem in wrapper.iter():
        if elem is wrapper:
            continue
        old_id = elem.get("id")
        if old_id:
            new_id = f"{prefix}-{old_id}"
            id_map[old_id] = new_id
            elem.set("id", new_id)
    if not id_map:
        return

    def _rewrite_value(value: str) -> str:
        new = value
        for old_id, new_id in id_map.items():
            new = new.replace(f"url(#{old_id})", f"url(#{new_id})")
            new = new.replace(f'url("#{old_id}")', f'url("#{new_id}")')
            new = new.replace(f"url('#{old_id}')", f"url('#{new_id}')")
        return new

    # 2) Rewrite reference attributes & style values.
    style_tag = f"{SVG_TAG}style"
    for elem in wrapper.iter():
        for attr in _REF_ATTRS:
            val = elem.get(attr)
            if not val:
                continue
            # Bare `#fragment` is an in-document reference; rewrite as a whole.
            if val.startswith("#"):
                old_id = val[1:]
                if old_id in id_map:
                    elem.set(attr, f"#{id_map[old_id]}")
                continue
            if "url(" in val:
                elem.set(attr, _rewrite_value(val))
        style = elem.get("style")
        if style and "url(" in style:
            elem.set("style", _rewrite_value(style))
        # `<style>…</style>` text nodes also carry CSS rules with url(#…)
        # references that many SVG exporters emit (Inkscape, Affinity, etc.).
        # Without rewriting these, an SVG with stylesheet-bound clip-paths
        # gets cross-bound to the wrong copy when the same source is reused.
        if elem.tag == style_tag and elem.text and "url(" in elem.text:
            elem.text = _rewrite_value(elem.text)


def fill_placeholders(template_path: Path, manifest: dict, output_path: Path) -> None:
    """Write a filled SVG with element bodies + label text + attribution caption."""
    parser = etree.XMLParser(remove_blank_text=False)
    tree = etree.parse(str(template_path), parser)
    root = tree.getroot()

    library_root = Path(
        manifest.get("library_root", "~/sci-illustration-library")
    ).expanduser()

    used_attribution: list[dict] = []

    # 1. Element placeholders.
    for panel_id, element_value in (manifest.get("panels") or {}).items():
        if _suspicious_id(panel_id):
            print(f"Warning: rejecting suspicious placeholder id: {panel_id!r}",
                  file=sys.stderr)
            continue

        element_full_path = resolve_element_path(element_value, library_root)
        if not element_full_path.exists():
            print(f"Warning: element not found: {element_full_path} "
                  f"(value='{element_value}')", file=sys.stderr)
            continue

        attr_record = lookup_attribution(element_value, library_root)
        # Track whether anything actually got inserted onto the canvas; we
        # only credit attribution after a successful fill so a failed parse
        # / missing placeholder / continue branch can't pollute the caption.
        inserted_any = False

        placeholders = _find_placeholders(root, panel_id)
        if not placeholders:
            print(f"Warning: placeholder '{panel_id}' not in template", file=sys.stderr)
            continue

        if element_full_path.suffix.lower() == ".svg":
            try:
                elem_tree = etree.parse(str(element_full_path))
                elem_root = elem_tree.getroot()
            except Exception as exc:
                print(f"Warning: failed to parse {element_full_path}: {exc}",
                      file=sys.stderr)
                continue

            # `_find_placeholders` may return multiple hits for one panel_id;
            # suffix the wrapper id with the 1-based index so the output SVG
            # never has duplicate ids (which break selectors / a11y / editor
            # validation). Inserting at the placeholder's original index
            # preserves SVG document order = z-order.
            for idx, placeholder in enumerate(placeholders, start=1):
                parent = placeholder.getparent()
                if parent is None:
                    continue
                wrapper_id = f"panel-{panel_id}-{idx}"
                wrapper = etree.Element(
                    f"{SVG_TAG}g",
                    attrib={
                        "id": wrapper_id,
                        "transform": placeholder.get("transform", ""),
                    },
                )
                # Deep-copy each child so two placeholders sharing one source
                # SVG don't share node identity.
                for child in elem_root:
                    wrapper.append(copy.deepcopy(child))
                # Namespace inner ids+refs so multiple copies of the same
                # source SVG don't produce id collisions or cross-bind
                # `url(#…)` references to the wrong instance.
                _namespace_subtree_ids(wrapper, wrapper_id)
                insert_at = list(parent).index(placeholder)
                parent.remove(placeholder)
                parent.insert(insert_at, wrapper)
                inserted_any = True
        else:
            # Use a file:// URI rather than a bare filesystem path so the
            # href stays valid across spaces in path components, Windows
            # drive letters, and `xmllint`/browser SVG renderers that treat
            # naked paths as relative.
            href_uri = element_full_path.resolve().as_uri()
            for idx, placeholder in enumerate(placeholders, start=1):
                parent = placeholder.getparent()
                if parent is None:
                    continue
                # Carry placeholder's `transform` so PNG/JPG/EMF land in the
                # same spot the SVG branch would: scale / rotate / translate
                # set on the placeholder all apply to the resulting <image>.
                image_attrib = {
                    "id": f"panel-{panel_id}-{idx}",
                    "x": placeholder.get("x", "0"),
                    "y": placeholder.get("y", "0"),
                    "width": placeholder.get("width", "100"),
                    "height": placeholder.get("height", "100"),
                    XLINK_HREF: href_uri,
                }
                ph_transform = placeholder.get("transform")
                if ph_transform:
                    image_attrib["transform"] = ph_transform
                image_elem = etree.Element(f"{SVG_TAG}image", attrib=image_attrib)
                insert_at = list(parent).index(placeholder)
                parent.remove(placeholder)
                parent.insert(insert_at, image_elem)
                inserted_any = True

        # Only count this asset toward the figure's attribution once at
        # least one placeholder was actually filled in.
        if inserted_any and attr_record is not None:
            used_attribution.append(attr_record)

        print(f"[fill] panel {panel_id} <- {element_value}")

    # 2. Text labels.
    for label_id, text_content in (manifest.get("labels") or {}).items():
        if _suspicious_id(label_id):
            print(f"Warning: rejecting suspicious label id: {label_id!r}", file=sys.stderr)
            continue
        text_elems = root.findall(
            f'.//svg:text[@id="{label_id}"]',
            namespaces=SVG_NS,
        )
        if not text_elems:
            print(f"Warning: label '{label_id}' not in template", file=sys.stderr)
            continue
        for t in text_elems:
            for child in list(t):
                t.remove(child)
            t.text = str(text_content)
            print(f"[fill] label {label_id} <- '{text_content}'")

    # 3. Attribution caption.
    if used_attribution:
        attr_lines = sorted({
            r.get("attribution") or f"{r.get('source_name', '')} ({r.get('license', '')})"
            for r in used_attribution
        })
        attr_text = "Adapted from: " + "; ".join(attr_lines)
        # Per SVG spec, viewBox values are separated by whitespace, comma, or
        # both (e.g. "0, 0, 1600, 1000" is valid). `.split()` only handles
        # whitespace and would either silently fall back to defaults or, worse,
        # raise ValueError on `float("1600,")`. Use a comma+whitespace regex.
        viewbox_raw = (root.get("viewBox") or "0 0 1600 1000").strip()
        viewbox = [v for v in re.split(r"[\s,]+", viewbox_raw) if v]
        # viewBox = "min-x min-y width height". Bottom-right is at
        # (min-x + width, min-y + height); not (width, height) when origin
        # is non-zero. Templates with a translated origin (e.g.
        # `viewBox="100 200 1600 1000"`) would otherwise place the
        # caption inside the visible area but offset from the corner.
        origin_x = float(viewbox[0]) if len(viewbox) >= 4 else 0.0
        origin_y = float(viewbox[1]) if len(viewbox) >= 4 else 0.0
        canvas_w = float(viewbox[2]) if len(viewbox) >= 4 else 1600.0
        canvas_h = float(viewbox[3]) if len(viewbox) >= 4 else 1000.0
        attr_node = etree.SubElement(
            root, f"{SVG_TAG}text",
            attrib={
                "x": str(origin_x + canvas_w - 10),
                "y": str(origin_y + canvas_h - 5),
                "text-anchor": "end",
                "font-size": "8",
                "fill": "#888",
                "font-family": "Arial",
                "class": "attribution",
            },
        )
        attr_node.text = attr_text
        print(f"[attribution] injected: {attr_text[:80]}...")

    tree.write(str(output_path), xml_declaration=True,
               encoding="utf-8", pretty_print=True)
    print(f"[fill] saved: {output_path}")


# ---------------------------------------------------------------------------
# Font normalization & exports
# ---------------------------------------------------------------------------


def normalize_fonts(svg_path: Path, font_family: str) -> None:
    """Force every text-bearing element to use one font-family.

    Three passes:
      1. set font-family on the root <svg> so inheritance picks it up;
      2. rewrite any explicit `font-family=...` attribute or `style=...;font-family:...`;
      3. for `<text>` / `<tspan>` that declare neither, set an explicit
         `font-family` so merged CC0 subtrees (whose own <svg> root we
         dropped) don't fall back to the UA default.
    """
    tree = etree.parse(str(svg_path))
    root = tree.getroot()

    count = 0
    # (1) root inheritance.
    if root.get("font-family") != font_family:
        root.set("font-family", font_family)
        count += 1

    text_tags = {f"{SVG_TAG}text", f"{SVG_TAG}tspan"}
    for elem in root.iter():
        # (2) rewrite explicit attribute.
        if elem.get("font-family") and elem.get("font-family") != font_family:
            elem.set("font-family", font_family)
            count += 1

        # (2) rewrite within style="...".
        style = elem.get("style", "")
        if "font-family" in style:
            parts: list[str] = []
            for chunk in style.split(";"):
                if chunk.strip().startswith("font-family"):
                    parts.append(f"font-family:{font_family}")
                else:
                    parts.append(chunk)
            elem.set("style", ";".join(parts))
            count += 1

        # (3) explicitly set on text-bearing elements that have nothing.
        if (elem.tag in text_tags
                and not elem.get("font-family")
                and "font-family" not in (elem.get("style") or "")):
            elem.set("font-family", font_family)
            count += 1

    tree.write(str(svg_path), xml_declaration=True, encoding="utf-8")
    print(f"[font] normalized {count} elements -> {font_family}")


def _run_inkscape(cmd: list[str], purpose: str) -> bool:
    """Run an Inkscape CLI invocation; return False (don't crash) on failure."""
    try:
        subprocess.run(cmd, check=True, capture_output=True)
        return True
    except FileNotFoundError:
        print(f"Warning: Inkscape CLI not found on PATH; skipping {purpose}.",
              file=sys.stderr)
        print("  Install Inkscape >= 1.2 from https://inkscape.org/", file=sys.stderr)
        return False
    except subprocess.CalledProcessError as exc:
        stderr_tail = (exc.stderr or b"").decode("utf-8", errors="replace")[-400:]
        print(f"Warning: Inkscape failed during {purpose}: {exc.returncode}",
              file=sys.stderr)
        if stderr_tail.strip():
            print(f"  stderr: {stderr_tail.strip()}", file=sys.stderr)
        return False


def export_pdf(svg_path: Path, pdf_path: Path, outline_text: bool = False) -> None:
    """Export PDF via Inkscape; for submission, set outline_text=True."""
    cmd = [
        "inkscape", str(svg_path),
        "--export-type=pdf",
        "--export-pdf-version=1.5",
        f"--export-text-to-path={'true' if outline_text else 'false'}",
        f"--export-filename={pdf_path}",
    ]
    if _run_inkscape(cmd, f"PDF export -> {pdf_path}"):
        print(f"[export] PDF: {pdf_path}")


def export_pptx(svg_path: Path, pptx_path: Path) -> None:
    """Render the SVG to a high-DPI PNG and embed into a 16:9 PPTX slide."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("Warning: python-pptx not installed; skipping PPTX export.",
              file=sys.stderr)
        print("  pip install python-pptx", file=sys.stderr)
        return

    png_temp = svg_path.with_suffix(".tmp.png")
    rendered = _run_inkscape(
        [
            "inkscape", str(svg_path),
            "--export-type=png",
            "--export-dpi=300",
            f"--export-filename={png_temp}",
        ],
        "PNG raster for PPTX",
    )
    if not rendered:
        return

    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(png_temp),
            left=Inches(0.5), top=Inches(0.5),
            width=Inches(12.333), height=Inches(6.5),
        )
        prs.save(str(pptx_path))
        print(f"[export] PPTX: {pptx_path}")
    finally:
        png_temp.unlink(missing_ok=True)


# ---------------------------------------------------------------------------
# Per-journal validation
# ---------------------------------------------------------------------------


def _journal_config_dir(start: Path) -> Path:
    """Locate `_journal-configs/` by walking up from the script's location.

    Returns the first existing `_journal-configs/` directory found by walking
    parents up to the repo root, or a sensible default if none exists.
    """
    cur = start.resolve()
    for _ in range(6):
        candidate = cur / "_journal-configs"
        if candidate.is_dir():
            return candidate
        if cur.parent == cur:
            break
        cur = cur.parent
    return start / "_journal-configs"  # default; load_journal_config will error


def load_journal_config(name: str, search_root: Path | None = None) -> dict:
    """Load `<name>.yaml` from `_journal-configs/` and validate the schema."""
    cfg_dir = _journal_config_dir(search_root or Path(__file__).parent)
    cfg_path = cfg_dir / f"{name}.yaml"
    if not cfg_path.exists():
        available = sorted(p.stem for p in cfg_dir.glob("*.yaml") if p.stem != "README")
        raise FileNotFoundError(
            f"Journal config not found: {cfg_path}\n"
            f"Available: {', '.join(available) if available else '(none)'}"
        )
    cfg = yaml.safe_load(cfg_path.read_text()) or {}
    if not isinstance(cfg, dict) or "name" not in cfg:
        raise ValueError(f"{cfg_path} is not a valid journal config")
    return cfg


def _iter_text_elements(root: etree._Element):
    """Yield every `<text>` and `<tspan>` descendant."""
    text_tags = {f"{SVG_TAG}text", f"{SVG_TAG}tspan"}
    for elem in root.iter():
        if elem.tag in text_tags:
            yield elem


def _get_inherited_attr(elem: etree._Element, attr: str) -> str | None:
    """Walk the parent chain looking for an attribute."""
    cur: etree._Element | None = elem
    while cur is not None:
        val = cur.get(attr)
        if val:
            return val
        cur = cur.getparent()
    return None


def _parse_pt(value: str | None) -> float | None:
    """Convert a CSS-ish font-size to pt (1pt ≈ 1.333 px)."""
    if not value:
        return None
    s = value.strip().lower()
    try:
        if s.endswith("pt"):
            return float(s[:-2])
        if s.endswith("px"):
            return float(s[:-2]) * 0.75  # px → pt
        if s.endswith("mm"):
            return float(s[:-2]) * 2.834645669
        if s.endswith("pc"):
            return float(s[:-2]) * 12.0
        # bare number — treat as user units (~px in our templates)
        return float(s) * 0.75
    except ValueError:
        return None


def validate_against_journal(svg_path: Path, journal_cfg: dict) -> tuple[list[str], list[str]]:
    """Return (errors, warnings) — errors fail the build, warnings just print.

    Checks:
      - fonts.forbidden present anywhere → error
      - fonts.preferred not used anywhere → warning
      - font_size.min_pt / max_pt violated → error / warning
      - line_weight.min_pt violated → error
      - sizes.* exceeded → warning
    """
    errors: list[str] = []
    warnings: list[str] = []

    try:
        tree = etree.parse(str(svg_path))
    except (etree.XMLSyntaxError, etree.ParseError, OSError) as exc:
        return ([f"could not parse SVG: {exc}"], [])
    root = tree.getroot()

    fonts_cfg = journal_cfg.get("fonts") or {}
    forbidden = {f.lower() for f in fonts_cfg.get("forbidden", [])}
    preferred = {f.lower() for f in fonts_cfg.get("preferred", [])}
    fs_cfg = journal_cfg.get("font_size") or {}
    min_pt = float(fs_cfg.get("min_pt", 0))
    max_pt = float(fs_cfg.get("max_pt", 9999))
    lw_cfg = journal_cfg.get("line_weight") or {}
    min_lw_pt = float(lw_cfg.get("min_pt", 0))

    used_fonts: set[str] = set()

    # Text-bearing elements: font + size checks.
    for elem in _iter_text_elements(root):
        ff = (_get_inherited_attr(elem, "font-family") or "").strip().strip('"\'')
        if ff:
            for chunk in ff.split(","):
                used_fonts.add(chunk.strip().strip('"\'').lower())
        fs_raw = _get_inherited_attr(elem, "font-size")
        fs_pt = _parse_pt(fs_raw)
        if fs_pt is not None:
            if fs_pt < min_pt:
                errors.append(f"<{elem.tag.split('}')[-1]} id={elem.get('id')!r}>"
                              f" font-size {fs_pt:.1f}pt < {min_pt}pt min")
            if fs_pt > max_pt:
                warnings.append(f"<{elem.tag.split('}')[-1]} id={elem.get('id')!r}>"
                                f" font-size {fs_pt:.1f}pt > {max_pt}pt max")

    # Substring match so the config can list canonical names like
    # "Comic Sans" while the SVG carries the OS-specific variant
    # "Comic Sans MS". Same logic for preferred matching.
    forbidden_hits = {
        bad for bad in forbidden
        if any(bad in used for used in used_fonts)
    }
    if forbidden_hits:
        errors.append(f"forbidden fonts present: {sorted(forbidden_hits)} "
                      f"(matched in: {sorted(used_fonts)})")
    if preferred:
        any_preferred_used = any(
            pref in used for pref in preferred for used in used_fonts
        )
        if not any_preferred_used:
            warnings.append(f"none of the preferred fonts {sorted(preferred)} "
                            f"were used (found: {sorted(used_fonts)})")

    # Stroke widths.
    for elem in root.iter():
        sw_raw = elem.get("stroke-width")
        sw_pt = _parse_pt(sw_raw)
        if sw_pt is not None and 0 < sw_pt < min_lw_pt:
            tag = elem.tag.split("}")[-1]
            errors.append(f"<{tag} id={elem.get('id')!r}> "
                          f"stroke-width {sw_pt:.2f}pt < {min_lw_pt}pt min")

    return errors, warnings


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------


def main() -> int:
    parser = argparse.ArgumentParser(description="Template-driven scientific figure assembler")
    parser.add_argument("--template", type=Path, required=True, help="SVG template")
    parser.add_argument("--manifest", type=Path, required=True, help="manifest YAML")
    parser.add_argument("--output", type=Path, required=True, help="output SVG path")
    parser.add_argument("--font", default="Arial", help="unified font-family (default: Arial)")
    parser.add_argument("--journal", type=str, default=None,
                        help=("validate the output against `_journal-configs/<name>.yaml` "
                              "(e.g. 'nature', 'cell', 'lancet'). Build fails on any error."))
    parser.add_argument("--export-pdf", action="store_true",
                        help="export an editable PDF alongside the SVG")
    parser.add_argument("--export-pdf-final", action="store_true",
                        help="export submission PDF with text converted to paths")
    parser.add_argument("--export-pptx", action="store_true",
                        help="export a 16:9 PPTX with the rendered figure embedded")
    args = parser.parse_args()

    if not args.template.exists():
        print(f"Error: template not found: {args.template}", file=sys.stderr)
        return 1
    if not args.manifest.exists():
        print(f"Error: manifest not found: {args.manifest}", file=sys.stderr)
        return 1

    args.output.parent.mkdir(parents=True, exist_ok=True)
    manifest = load_manifest(args.manifest)

    fill_placeholders(args.template, manifest, args.output)
    normalize_fonts(args.output, args.font)

    # Run journal validation BEFORE the (more expensive) export step so the
    # user sees violations early. Errors fail the build with non-zero exit.
    if args.journal:
        try:
            cfg = load_journal_config(args.journal)
        except (FileNotFoundError, ValueError) as exc:
            print(f"Error: {exc}", file=sys.stderr)
            return 1
        errors, warnings = validate_against_journal(args.output, cfg)
        if warnings:
            print(f"\n[journal:{cfg['name']}] warnings ({len(warnings)}):", file=sys.stderr)
            for w in warnings:
                print(f"  - {w}", file=sys.stderr)
        if errors:
            print(f"\n[journal:{cfg['name']}] errors ({len(errors)}):", file=sys.stderr)
            for e in errors:
                print(f"  - {e}", file=sys.stderr)
            print("\nFix the errors above or remove --journal to bypass.",
                  file=sys.stderr)
            return 1
        if not warnings and not errors:
            print(f"\n[journal:{cfg['name']}] ✓ all checks passed")

    if args.export_pdf:
        export_pdf(
            args.output,
            args.output.with_name(args.output.stem + "_editable.pdf"),
            outline_text=False,
        )
    if args.export_pdf_final:
        export_pdf(
            args.output,
            args.output.with_name(args.output.stem + "_submission.pdf"),
            outline_text=True,
        )
    if args.export_pptx:
        export_pptx(args.output, args.output.with_suffix(".pptx"))

    print(f"\nAll done. Main output: {args.output}")
    return 0


if __name__ == "__main__":
    sys.exit(main())

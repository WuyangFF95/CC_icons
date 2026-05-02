#!/usr/bin/env python3
"""
library_tools.py — manage the local scientific element library.

Commands:
  qc          — quality-check a single SVG element (no registration side effect)
  register    — register a QC-passed element into a per-category _index.yaml
  search      — full-text search across the unified index (CC0 bulk + legacy YAML)
  audit       — sweep every category and list every element that fails QC
  check       — given a project manifest, list which required elements are
                already in the library and which are missing
  stats       — summary of the unified index by source / category / license
  preview     — render a thumbnail grid PNG of the matched elements
  export      — bundle matched elements into a one-element-per-slide PPTX
  attribution — emit a markdown attribution block for CC-BY elements
                (use before submitting a figure)

Dependencies:
  - lxml, pyyaml         (always)
  - cairosvg, pillow     (preview, export)
  - python-pptx          (export)

Usage:
  python library_tools.py qc cells/treg-flat-v1.svg
  python library_tools.py register cells/treg-flat-v1.svg \
         --subject "regulatory T cell" --style flat-blue --provider recraft-v3
  python library_tools.py search "T cell" --license CC0
  python library_tools.py stats
  python library_tools.py preview --query "neuron" --output preview.png
"""

from __future__ import annotations

import argparse
import json
import re
import sys
import tempfile
from collections import Counter
from datetime import date
from io import BytesIO
from pathlib import Path

import yaml
from lxml import etree

# Optional deps — lazy-checked per command.
# Catch both ImportError (package missing) and OSError (e.g. cairosvg installed
# but libcairo system library missing on the host).
try:
    import cairosvg
    from PIL import Image, ImageDraw, ImageFont
    HAS_PREVIEW = True
except (ImportError, OSError):
    HAS_PREVIEW = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    HAS_PPTX = True
except (ImportError, OSError):
    HAS_PPTX = False


SVG_NS = {"svg": "http://www.w3.org/2000/svg"}
# Mutated in main() when --library-root is supplied.
LIBRARY_ROOT = Path.home() / "sci-illustration-library"
# Top-level legacy categories. The unified index also uses sub-categories
# like "cells/immune" and "equipment/lab"; CLI filters accept either form.
CATEGORIES = [
    "cells", "molecules", "organelles", "tissues",
    "organs", "equipment", "pathways", "arrows",
]


# ---------------------------------------------------------------------------
# Unified index loader (CC0 bulk JSON + legacy per-category YAML)
# ---------------------------------------------------------------------------


def load_unified_index() -> list[dict]:
    """Merge the CC0 bulk index.json with each category's legacy _index.yaml.

    Returns a list of records sharing one schema; legacy records carry
    `_legacy: True` and a `_meta` field with the original YAML row.
    """
    records: list[dict] = []
    seen_ids: set[str] = set()

    # 1. CC0 bulk index produced by download_cc0_seed.py.
    json_index = LIBRARY_ROOT / "library" / "index.json"
    if json_index.exists():
        try:
            data = json.loads(json_index.read_text())
            for r in data:
                rid = r.get("id")
                if rid and rid not in seen_ids:
                    seen_ids.add(rid)
                    r.setdefault("_root", json_index.parent)
                    records.append(r)
        except Exception as exc:
            print(f"[warn] failed to parse {json_index}: {exc}", file=sys.stderr)

    # 2. Legacy per-category YAML (manually registered elements).
    for category in CATEGORIES:
        yaml_index = LIBRARY_ROOT / category / "_index.yaml"
        if not yaml_index.exists():
            continue
        try:
            idx = yaml.safe_load(yaml_index.read_text()) or {"elements": []}
            for entry in idx.get("elements", []):
                fid = f"manual-{category}-{entry.get('file', '').replace('.svg', '')}"
                if fid in seen_ids:
                    continue
                seen_ids.add(fid)
                records.append({
                    "id": fid,
                    "name": entry.get("subject", ""),
                    "tags": [entry.get("style", "")],
                    "category": category,
                    "source": entry.get("provider", "manual"),
                    "license": entry.get("license", "research-use-only"),
                    "license_url": "",
                    "attribution_required": False,
                    "file": str(Path(category) / entry.get("file", "")),
                    "_root": LIBRARY_ROOT,
                    "_legacy": True,
                    "_meta": entry,
                })
        except Exception as exc:
            print(f"[warn] failed to parse {yaml_index}: {exc}", file=sys.stderr)

    return records


def resolve_svg_path(record: dict) -> Path:
    """Convert a unified-index record into an absolute SVG path."""
    root = record.get("_root", LIBRARY_ROOT / "library")
    return Path(root) / record["file"]


# ---------------------------------------------------------------------------
# Quality check
# ---------------------------------------------------------------------------


def quality_check(svg_path: Path) -> dict:
    """Inspect one SVG for path count, color palette, embedded text, etc."""
    try:
        tree = etree.parse(str(svg_path))
        root = tree.getroot()
    except Exception as exc:
        return {"pass": False, "error": str(exc)}

    paths = root.findall(".//svg:path", SVG_NS)
    texts = root.findall(".//svg:text", SVG_NS)
    images = root.findall(".//svg:image", SVG_NS)
    gradients = (root.findall(".//svg:linearGradient", SVG_NS)
                 + root.findall(".//svg:radialGradient", SVG_NS))

    total_nodes = sum(
        len(re.findall(r"[MLCQAHVZmlcqahvz]", p.get("d", "")))
        for p in paths
    )

    fills: set[str] = set()
    for elem in root.iter():
        fill = elem.get("fill") or ""
        if (fill and fill not in ("none", "transparent")
                and not fill.startswith("url")):
            fills.add(fill.lower())

    checks = {
        "path_count":       (len(paths) < 200, f"{len(paths)} paths"),
        "no_text_nodes":    (len(texts) == 0, f"{len(texts)} <text> nodes"),
        "no_image_nodes":   (len(images) == 0, f"{len(images)} <image> nodes"),
        "node_count":       (total_nodes < 2000, f"{total_nodes} path commands"),
        "color_palette":    (len(fills) <= 8, f"{len(fills)} unique colors"),
        "gradient_minimal": (len(gradients) <= 2, f"{len(gradients)} gradients"),
    }

    return {
        "pass": all(passed for passed, _ in checks.values()),
        "checks": checks,
        "stats": {
            "paths": len(paths),
            "nodes": total_nodes,
            "colors": len(fills),
            "gradients": len(gradients),
            "texts": len(texts),
            "images": len(images),
        },
    }


# ---------------------------------------------------------------------------
# Commands
# ---------------------------------------------------------------------------


def cmd_qc(args: argparse.Namespace) -> int:
    """Print a QC report for one SVG."""
    path = Path(args.file)
    if not path.exists():
        print(f"Error: file not found: {path}", file=sys.stderr)
        return 1

    result = quality_check(path)
    print(f"\n=== QC Report: {path.name} ===\n")

    if "error" in result:
        print(f"FATAL: {result['error']}")
        return 1

    for check_name, (passed, detail) in result["checks"].items():
        symbol = "PASS" if passed else "FAIL"
        print(f"  [{symbol}] {check_name:20s}: {detail}")

    verdict = "ACCEPT (ready to register)" if result["pass"] else "REJECT (regenerate or simplify)"
    print(f"\nVerdict: {verdict}\n")
    return 0 if result["pass"] else 1


def get_category_for_subject(subject: str) -> str:
    """Best-effort mapping from a subject string to a top-level category."""
    subject_lower = subject.lower()
    rules = {
        "cells": ["cell", "cyte", "phage", "blast"],
        "molecules": ["protein", "antibody", "dna", "rna", "enzyme",
                      "receptor", "ligand"],
        "organelles": ["mitochond", "nucleus", "golgi", "endoplasmic",
                       "lysosome", "ribosome"],
        "tissues": ["tissue", "epithelium", "stroma", "matrix",
                    "glomerul", "lobule"],
        "organs": ["liver", "kidney", "lung", "brain", "heart",
                   "pancreas", "spleen"],
        "equipment": ["microscope", "cytometer", "pcr", "plate",
                      "tube", "chamber", "array"],
        "pathways": ["pathway", "cascade", "signaling", "phospho",
                     "transport"],
    }
    for category, keywords in rules.items():
        if any(kw in subject_lower for kw in keywords):
            return category
    return "molecules"  # fallback


def cmd_register(args: argparse.Namespace) -> int:
    """Register a QC-passed SVG into the matching category's _index.yaml."""
    src = Path(args.file)
    if not src.exists():
        print(f"Error: file not found: {src}", file=sys.stderr)
        return 1

    qc = quality_check(src)
    if not qc["pass"]:
        print("QC FAILED. Cannot register. Run `qc` for details.", file=sys.stderr)
        return 1

    category = args.category or get_category_for_subject(args.subject)
    target_dir = LIBRARY_ROOT / category
    target_dir.mkdir(parents=True, exist_ok=True)

    safe_subject = re.sub(r"[^a-z0-9]+", "-", args.subject.lower()).strip("-")
    target_file = f"{safe_subject}-{args.style}-v{args.version}.svg"
    target_path = target_dir / target_file
    target_path.write_bytes(src.read_bytes())

    index_path = target_dir / "_index.yaml"
    if index_path.exists():
        index = yaml.safe_load(index_path.read_text()) or {"elements": []}
    else:
        index = {"elements": []}

    entry = {
        "file": target_file,
        "subject": args.subject,
        "style": args.style,
        "style_ref": args.style_ref or "",
        "provider": args.provider,
        "style_id": args.style_id or "",
        "generated_at": str(date.today()),
        "nodes": qc["stats"]["nodes"],
        "colors": qc["stats"]["colors"],
        "qc_passed": True,
        "used_in": [],
        "license": args.license,
    }
    # Upsert by file: re-running `register` overwrites the on-disk SVG, and
    # the metadata must reflect the new run. Otherwise `load_unified_index()`
    # dedupes by derived id and pins the *first* row, so search/export/
    # attribution would keep showing stale metadata.
    index["elements"] = [
        existing for existing in index["elements"]
        if existing.get("file") != target_file
    ]
    index["elements"].append(entry)
    index_path.write_text(yaml.safe_dump(index, allow_unicode=True, sort_keys=False))

    print(f"Registered: {target_path}")
    print(f"Updated index: {index_path}")
    return 0


def cmd_search(args: argparse.Namespace) -> int:
    """Full-text search the unified index."""
    records = load_unified_index()
    if not records:
        # Empty CSV (just a newline) when piping; human prompt otherwise.
        if not args.ids_only:
            print("Library is empty. Run download_cc0_seed.py to seed.")
        return 0

    query = (args.query or "").lower().strip()
    matches: list[dict] = []
    for r in records:
        haystack = " ".join([
            r.get("name", ""),
            r.get("category", ""),
            r.get("id", ""),
            " ".join(r.get("tags", [])),
        ]).lower()

        if query and query not in haystack:
            continue
        if args.category:
            cat = r.get("category", "")
            # `--category cells` matches "cells" plus any "cells/<sub>".
            # `--category cells/immune` matches that record exactly.
            if not (cat == args.category or cat.startswith(args.category + "/")):
                continue
        if args.source and args.source not in r.get("source", ""):
            continue
        if args.license and args.license.upper() not in r.get("license", "").upper():
            continue
        # Backward-compat flags from v0.1.0.
        if args.subject and args.subject.lower() not in r.get("name", "").lower():
            continue
        if args.style and not any(args.style in t for t in r.get("tags", [])):
            continue
        if args.provider and args.provider not in r.get("source", ""):
            continue

        matches.append(r)

    if args.max:
        matches = matches[: args.max]

    if not matches:
        if not args.ids_only:
            print("No matches.")
        return 0

    if args.ids_only:
        # Pure CSV on stdout, nothing else — pipe-safe (e.g. `--ids "$(... --ids-only)"`).
        print(",".join(r["id"] for r in matches))
        return 0

    print(f"\nFound {len(matches)} element(s):\n")
    for r in matches:
        attr_flag = " [ATTR]" if r.get("attribution_required") else ""
        print(f"  {r['id']}{attr_flag}")
        print(f"    name:     {r['name']}")
        print(f"    category: {r['category']}")
        print(f"    source:   {r['source']}  ({r['license']})")
        print(f"    file:     {r['file']}")
        print()

    return 0


def cmd_audit(args: argparse.Namespace) -> int:
    """QC every SVG in the unified index (CC0 bulk + legacy YAML rows).

    Walking the unified index instead of a CATEGORIES dir glob means
    bootstrapped CC0 assets under e.g. `library/cells/immune/` are no longer
    silently skipped — those rows are merged into the index by
    `load_unified_index()`. Non-SVG assets (.png/.ai/.emf) are out of scope
    for QC (it inspects path/text/gradient counts) and are skipped here.
    """
    failed: list[tuple[Path, dict]] = []
    total = 0
    seen: set[Path] = set()
    for record in load_unified_index():
        svg = resolve_svg_path(record)
        if svg.suffix.lower() != ".svg":
            continue
        try:
            resolved = svg.resolve()
        except OSError:
            continue
        if resolved in seen or not svg.exists():
            continue
        seen.add(resolved)
        total += 1
        qc = quality_check(svg)
        if not qc["pass"]:
            failed.append((svg, qc))

    print(f"\nAudit complete: {total} elements scanned, {len(failed)} failed QC.\n")
    for svg, qc in failed:
        # `try/except` guards against unified-index paths outside LIBRARY_ROOT
        # (would happen if a record's `_root` points elsewhere).
        try:
            label = svg.relative_to(LIBRARY_ROOT)
        except ValueError:
            label = svg
        print(f"  FAIL: {label}")
        for check_name, (passed, detail) in qc["checks"].items():
            if not passed:
                print(f"    - {check_name}: {detail}")
        print()
    return 0 if not failed else 1


def cmd_check(args: argparse.Namespace) -> int:
    """Compare a project manifest's element list against what's registered."""
    manifest_path = LIBRARY_ROOT / "_final" / args.project / "manifest.yaml"
    if not manifest_path.exists():
        print(f"Error: manifest not found: {manifest_path}", file=sys.stderr)
        print("Hint: create one with project's element requirements.")
        return 1

    manifest = yaml.safe_load(manifest_path.read_text()) or {}
    print(f"\n=== Project: {args.project} ===\n")
    have: list[tuple[dict, str, str]] = []
    missing: list[dict] = []

    # Walk the unified index once so v0.1.1 CC0 records show up as "have"
    # alongside legacy YAML rows. The legacy walk stays as the second tier
    # to preserve the v0.1.0 contract of matching on `subject`+`style`.
    unified_records = load_unified_index()

    for req in manifest.get("elements_needed", []):
        found = False
        # Tier 1: unified index (CC0 bulk + already-merged legacy rows).
        # CC0 record `tags` carry content categories (e.g. ["immune"],
        # ["silhouette"]), not visual style names like "flat-blue", so the
        # style guard MUST treat an unspecified style as a wildcard rather
        # than as the empty string "in" tags. Also, when a style IS
        # specified, an empty `tags` list still shouldn't auto-fail — fall
        # through to the legacy YAML tier below.
        req_style = (req.get("style") or "").lower()
        for r in unified_records:
            name = (r.get("name") or "").lower()
            tags = [(t or "").lower() for t in r.get("tags", [])]
            style_ok = (not req_style) or (req_style in tags)
            if req["subject"].lower() in name and style_ok:
                have.append((req, r.get("category", ""), r.get("file", "")))
                found = True
                break

        # Tier 2: per-category legacy YAML (preserves v0.1.0 contract).
        if not found:
            for category in CATEGORIES:
                idx_path = LIBRARY_ROOT / category / "_index.yaml"
                if not idx_path.exists():
                    continue
                idx = yaml.safe_load(idx_path.read_text()) or {"elements": []}
                for entry in idx.get("elements", []):
                    if (req["subject"].lower() in entry.get("subject", "").lower()
                            and entry.get("style") == req.get("style")):
                        have.append((req, category, entry["file"]))
                        found = True
                        break
                if found:
                    break

        if not found:
            missing.append(req)

    print(f"In library ({len(have)}):")
    for req, category, file in have:
        print(f"  [OK] {req['subject']:30s} -> {category}/{file}")

    print(f"\nMissing ({len(missing)}):")
    for req in missing:
        print(f"  [..] {req['subject']:30s} (style: {req.get('style')})")

    if missing:
        print(f"\nNext step: generate {len(missing)} missing elements via "
              "recraft-scientific-illustration skill.")
    return 0


# ---------------------------------------------------------------------------
# New commands: stats / preview / export / attribution
# ---------------------------------------------------------------------------


def cmd_stats(args: argparse.Namespace) -> int:
    """Print a one-shot summary of the unified library."""
    records = load_unified_index()
    if not records:
        print("Library is empty.")
        return 0

    by_source = Counter(r.get("source", "?") for r in records)
    by_category = Counter(r.get("category", "?") for r in records)
    by_license = Counter(r.get("license", "?") for r in records)

    print(f"\n{'=' * 60}")
    print(f"Library Stats  (root: {LIBRARY_ROOT})")
    print(f"{'=' * 60}")
    print(f"\nTotal elements: {len(records)}\n")

    print("By source:")
    for s, n in by_source.most_common():
        print(f"  {s:25s} {n:>6d}")

    print("\nBy category:")
    for c, n in by_category.most_common():
        print(f"  {c:25s} {n:>6d}")

    print("\nBy license:")
    for lic, n in by_license.most_common():
        print(f"  {lic:25s} {n:>6d}")

    attr_required = sum(1 for r in records if r.get("attribution_required"))
    print(f"\nAttribution required: {attr_required} elements")
    return 0


def _filter_for_grid(args: argparse.Namespace, records: list[dict]) -> list[dict]:
    """Shared selector for preview / export."""
    query = (args.query or "").lower()
    matches = [
        r for r in records
        if not query
        or query in r.get("name", "").lower()
        or query in r.get("category", "").lower()
        or query in " ".join(r.get("tags", [])).lower()
    ]
    if args.ids:
        wanted = set(args.ids.split(","))
        matches = [r for r in matches if r["id"] in wanted]
    return matches[: args.max]


def cmd_preview(args: argparse.Namespace) -> int:
    """Render a thumbnail grid of matched elements as a single PNG."""
    if not HAS_PREVIEW:
        print("Error: cmd preview requires `cairosvg` and `Pillow`", file=sys.stderr)
        print("       pip install cairosvg pillow", file=sys.stderr)
        return 1

    matches = _filter_for_grid(args, load_unified_index())
    if not matches:
        print("No matches to preview.")
        return 0

    cell = args.cell_size
    cols = args.cols
    rows = (len(matches) + cols - 1) // cols
    label_h = 28
    grid_w = cols * (cell + 8) + 8
    grid_h = rows * (cell + label_h + 12) + 8

    canvas = Image.new("RGB", (grid_w, grid_h), "white")
    draw = ImageDraw.Draw(canvas)
    try:
        font = ImageFont.truetype("DejaVuSans.ttf", 10)
    except OSError:
        font = ImageFont.load_default()

    print(f"Rendering {len(matches)} thumbnails ({cols}x{rows})...")
    for i, r in enumerate(matches):
        row, col = i // cols, i % cols
        x = 8 + col * (cell + 8)
        y = 8 + row * (cell + label_h + 12)

        svg_path = resolve_svg_path(r)
        if not svg_path.exists():
            draw.rectangle([x, y, x + cell, y + cell], outline="red")
            draw.text((x + 4, y + cell // 2), "missing", fill="red", font=font)
            continue

        ext = svg_path.suffix.lower()
        try:
            if ext == ".svg":
                png_bytes = cairosvg.svg2png(
                    url=str(svg_path),
                    output_width=cell,
                    output_height=cell,
                )
                img = Image.open(BytesIO(png_bytes)).convert("RGB")
            elif ext in (".png", ".jpg", ".jpeg"):
                img = Image.open(svg_path).convert("RGB")
                img.thumbnail((cell, cell))
            else:
                # AI / EPS / EMF / WMF — these need an external converter
                # (libreoffice, inkscape) and aren't worth a hard dep here.
                draw.rectangle([x, y, x + cell, y + cell], outline="#888")
                draw.text((x + 4, y + 4), f"unsupported {ext}", fill="#888", font=font)
                continue
            # Center-paste raster preview within the cell.
            paste_x = x + (cell - img.width) // 2
            paste_y = y + (cell - img.height) // 2
            canvas.paste(img, (paste_x, paste_y))
        except Exception as exc:
            draw.rectangle([x, y, x + cell, y + cell], outline="orange")
            draw.text((x + 4, y + 4),
                      f"render fail: {type(exc).__name__}",
                      fill="orange", font=font)

        name = r.get("name", "")[:30]
        draw.text((x, y + cell + 2), name, fill="black", font=font)
        draw.text(
            (x, y + cell + 14),
            f"{r['source']} | {r['license'][:10]}",
            fill="#666", font=font,
        )

    out = Path(args.output).expanduser()
    canvas.save(out)
    print(f"Saved preview: {out}")
    return 0


def cmd_export(args: argparse.Namespace) -> int:
    """Bundle matched elements into a one-element-per-slide PPTX."""
    if not HAS_PPTX:
        print("Error: cmd export requires `python-pptx`", file=sys.stderr)
        print("       pip install python-pptx", file=sys.stderr)
        return 1
    if not HAS_PREVIEW:
        print("Error: cmd export also needs `cairosvg` + `Pillow`", file=sys.stderr)
        return 1

    records = load_unified_index()
    by_id = {r["id"]: r for r in records}

    if args.ids:
        wanted_ids = args.ids.split(",")
        selected = [by_id[i] for i in wanted_ids if i in by_id]
    elif args.query:
        q = args.query.lower()
        selected = [r for r in records if q in r.get("name", "").lower()][: args.max]
    else:
        print("Error: --ids or --query required", file=sys.stderr)
        return 1

    if not selected:
        print("No elements to export.")
        return 0

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # TemporaryDirectory ensures the rasterized PNGs get cleaned up regardless
    # of which branch we exit through (success / cairosvg failure / ctrl-C).
    with tempfile.TemporaryDirectory(prefix="library_export_") as tmp_str:
        tmp_dir = Path(tmp_str)
        print(f"Exporting {len(selected)} elements to PPTX (tmp: {tmp_dir})...")
        for r in selected:
            svg_path = resolve_svg_path(r)
            if not svg_path.exists():
                continue
            ext = svg_path.suffix.lower()
            # python-pptx accepts .png/.jpg directly; .svg must be rasterized;
            # other formats (.ai/.eps/.emf/.wmf) need an external converter.
            if ext == ".svg":
                picture_path = tmp_dir / f"{r['id']}.png"
                try:
                    cairosvg.svg2png(
                        url=str(svg_path),
                        write_to=str(picture_path),
                        output_width=1200,
                    )
                except Exception as exc:
                    print(f"  [skip] {r['id']}: SVG render failed "
                          f"({type(exc).__name__}: {exc})",
                          file=sys.stderr)
                    continue
            elif ext in (".png", ".jpg", ".jpeg"):
                picture_path = svg_path
            else:
                print(f"  [skip] {r['id']}: unsupported format {ext} "
                      "(needs libreoffice or inkscape conversion)",
                      file=sys.stderr)
                continue

            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(str(picture_path), Inches(2.5), Inches(1.0),
                                     height=Inches(4.5))

            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3),
                                                 Inches(12), Inches(0.6))
            title_box.text_frame.text = r["name"]
            title_box.text_frame.paragraphs[0].font.size = Pt(24)
            title_box.text_frame.paragraphs[0].font.bold = True

            footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.6),
                                              Inches(12), Inches(0.6))
            footer.text_frame.text = (
                f"{r['id']}  |  source: {r['source']}  |  license: {r['license']}"
            )
            footer.text_frame.paragraphs[0].font.size = Pt(11)

            if r.get("attribution_required"):
                attr_p = footer.text_frame.add_paragraph()
                attr_p.text = f"ATTRIBUTION: {r.get('attribution', '')}"
                attr_p.font.size = Pt(10)

        out_path = Path(args.output).expanduser()
        prs.save(str(out_path))
    print(f"Saved: {out_path}")
    return 0


def cmd_attribution(args: argparse.Namespace) -> int:
    """Render a markdown attribution block for CC-BY elements."""
    records = load_unified_index()

    if args.ids:
        wanted = set(args.ids.split(","))
        # Even when the caller pipes in IDs explicitly, only emit attribution
        # for records that need it; otherwise feeding `search --ids-only` (CC0
        # included) into here would bloat the caption with public-domain rows.
        selected = [r for r in records
                    if r["id"] in wanted and r.get("attribution_required")]
    else:
        selected = [r for r in records if r.get("attribution_required")]

    if not selected:
        print("No attribution-required elements selected.")
        return 0

    grouped: dict[tuple[str, str], list[dict]] = {}
    for r in selected:
        key = (r.get("source", ""), r.get("license", ""))
        grouped.setdefault(key, []).append(r)

    lines = ["# Figure Asset Attributions\n"]
    lines.append(f"_Generated {date.today()}_\n")
    lines.append(f"This figure uses {len(selected)} non-CC0 assets requiring attribution.\n")

    for (source, lic), items in sorted(grouped.items()):
        lines.append(f"\n## {source} ({lic})\n")
        seen_attr: set[str] = set()
        for r in items:
            attr = r.get("attribution", "") or f"{r.get('source_name', '')} ({lic})"
            if attr in seen_attr:
                continue
            seen_attr.add(attr)
            lines.append(f"- {attr}")
        lines.append("")

    lines.append("\n---")
    lines.append("\n## Suggested figure caption text:\n")
    cited = sorted({(r.get("attribution") or r.get("source_name", "")) for r in selected})
    lines.append("> Schematic illustrations adapted from " + ", ".join(cited) + ".")

    output = "\n".join(lines)
    if args.output:
        Path(args.output).expanduser().write_text(output)
        print(f"Saved: {args.output}")
    else:
        print(output)
    return 0


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Local scientific element library management",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument("--library-root",
                        default=str(Path.home() / "sci-illustration-library"),
                        help=("override library root (default: "
                              "~/sci-illustration-library; matches the "
                              "--target flag of download_cc0_seed.py)"))
    sub = parser.add_subparsers(dest="command", required=True)

    p_qc = sub.add_parser("qc", help="QC a single SVG element")
    p_qc.add_argument("file", help="path to SVG")

    p_reg = sub.add_parser("register", help="register a QC-passed element")
    p_reg.add_argument("file", help="path to SVG")
    p_reg.add_argument("--subject", required=True,
                       help='element subject, e.g. "regulatory T cell"')
    p_reg.add_argument("--style", required=True,
                       help='style label, e.g. "flat-blue"')
    p_reg.add_argument("--provider", required=True,
                       choices=["recraft-v3", "gemini-imagen-4",
                                "zhipu-cogview-4", "minimax", "gpt-image",
                                "hand-drawn"])
    p_reg.add_argument("--version", default="1", help="version tag (default: 1)")
    p_reg.add_argument("--style-ref", help="style reference filename")
    p_reg.add_argument("--style-id", help="provider-specific style id")
    p_reg.add_argument("--category", choices=CATEGORIES,
                       help="category override (default: derived from subject)")
    p_reg.add_argument("--license", default="research-use-only")

    p_search = sub.add_parser("search", help="search the unified library index")
    p_search.add_argument("query", nargs="?", default=None,
                          help="full-text query (matches name/category/tags/id)")
    # Accept any category string — the unified index uses subcategories like
    # "cells/immune" and "equipment/lab" alongside the top-level CATEGORIES.
    p_search.add_argument("--category",
                          help=("restrict to category (top-level e.g. cells, "
                                "or subcategory e.g. cells/immune)"))
    p_search.add_argument("--source", help="restrict to source (bioicons/phylopic/...)")
    p_search.add_argument("--license", help="license keyword (CC0 / CC-BY / Public)")
    p_search.add_argument("--max", type=int, default=50, help="cap result count")
    p_search.add_argument("--ids-only", action="store_true",
                          help="print a comma-separated list of matched IDs (for piping)")
    # legacy filters from v0.1.0
    p_search.add_argument("--subject", help="[legacy] subject keyword")
    p_search.add_argument("--style", help="[legacy] style keyword")
    p_search.add_argument("--provider", help="[legacy] provider keyword")

    sub.add_parser("audit", help="QC every element in the library")

    p_check = sub.add_parser("check", help="compare a project manifest against the library")
    p_check.add_argument("--project", required=True, help="project directory under _final/")

    sub.add_parser("stats", help="summary of the unified index")

    p_prev = sub.add_parser("preview", help="thumbnail grid PNG of matched elements")
    p_prev.add_argument("--query", help="full-text query")
    p_prev.add_argument("--ids", help="comma-separated IDs")
    p_prev.add_argument("--max", type=int, default=64, help="max elements (default 64 = 8x8)")
    p_prev.add_argument("--cols", type=int, default=8)
    p_prev.add_argument("--cell-size", type=int, default=120, help="thumbnail edge px")
    p_prev.add_argument("--output", "-o", default="./preview.png")

    p_export = sub.add_parser("export", help="bundle elements into a PPTX")
    p_export.add_argument("--ids", help="comma-separated IDs")
    p_export.add_argument("--query", help="full-text query")
    p_export.add_argument("--max", type=int, default=20)
    p_export.add_argument("--output", "-o", default="./elements.pptx")

    p_attr = sub.add_parser("attribution",
                            help="emit markdown attribution block for CC-BY elements")
    p_attr.add_argument("--ids", help="comma-separated IDs (default: all attribution-required)")
    p_attr.add_argument("--output", "-o", help="output markdown path (default stdout)")

    return parser


def main() -> int:
    parser = build_parser()
    args = parser.parse_args()

    # Apply --library-root before any command runs. Mutating the module global
    # keeps the loader / resolver functions free of an extra plumbing arg.
    global LIBRARY_ROOT
    LIBRARY_ROOT = Path(args.library_root).expanduser().resolve()
    LIBRARY_ROOT.mkdir(parents=True, exist_ok=True)

    handlers = {
        "qc": cmd_qc,
        "register": cmd_register,
        "search": cmd_search,
        "audit": cmd_audit,
        "check": cmd_check,
        "stats": cmd_stats,
        "preview": cmd_preview,
        "export": cmd_export,
        "attribution": cmd_attribution,
    }
    return handlers[args.command](args)


if __name__ == "__main__":
    sys.exit(main())
